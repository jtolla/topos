import logging
import re
from pathlib import Path

from pydantic import BaseModel

logger = logging.getLogger(__name__)


class ExtractedDocument(BaseModel):
    """Result of content extraction from a file."""

    title: str
    text: str
    sections: list[dict] = []  # Optional, for future use


def extract_text_plain(path: str) -> ExtractedDocument:
    """Extract content from a plain text file."""
    encodings = ["utf-8", "latin-1", "cp1252"]

    for encoding in encodings:
        try:
            with open(path, encoding=encoding) as f:
                text = f.read()
            title = Path(path).stem
            return ExtractedDocument(title=title, text=text)
        except UnicodeDecodeError:
            continue

    raise ValueError(f"Could not decode file with any supported encoding: {path}")


def extract_pdf(path: str) -> ExtractedDocument:
    """Extract content from a PDF file using pdfminer.six."""
    from pdfminer.high_level import extract_text

    text = extract_text(path)
    title = Path(path).stem
    return ExtractedDocument(title=title, text=text)


def extract_docx(path: str) -> ExtractedDocument:
    """Extract content from a DOCX file using python-docx."""
    from docx import Document

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs]
    text = "\n".join(paragraphs)

    # Try to get title from document properties
    title = doc.core_properties.title or Path(path).stem

    return ExtractedDocument(title=title, text=text)


def extract_pptx(path: str) -> ExtractedDocument:
    """Extract content from a PPTX file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(path)
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_content))

    text = "\n\n".join(slides_text)
    title = Path(path).stem

    return ExtractedDocument(title=title, text=text)


# MIME type to extractor mapping
EXTRACTORS = {
    "text/plain": extract_text_plain,
    "application/pdf": extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
}

# File extension fallbacks
EXTENSION_MIME_MAP = {
    ".txt": "text/plain",
    ".md": "text/plain",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def extract_content(path: str, file_type: str) -> ExtractedDocument:
    """
    Extract text content from a file based on its type.

    Args:
        path: Path to the file
        file_type: MIME type of the file

    Returns:
        ExtractedDocument with title and text

    Raises:
        ValueError: If file type is not supported
        Exception: If extraction fails
    """
    # Try to find extractor by MIME type
    extractor = EXTRACTORS.get(file_type)

    # Fall back to extension-based detection
    if not extractor:
        ext = Path(path).suffix.lower()
        fallback_mime = EXTENSION_MIME_MAP.get(ext)
        if fallback_mime:
            extractor = EXTRACTORS.get(fallback_mime)

    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type} for file {path}")

    try:
        return extractor(path)
    except Exception as e:
        logger.exception(f"Failed to extract content from {path}: {e}")
        raise


class ChunkSpec(BaseModel):
    """Specification for a single chunk of text."""

    index: int
    text: str
    char_start: int
    char_end: int
    section_path: list[str] | None = None  # Hierarchical path for type-aware chunking


def normalize_whitespace(text: str) -> str:
    """Normalize whitespace in text."""
    # Replace multiple whitespace with single space
    text = re.sub(r"[ \t]+", " ", text)
    # Replace multiple newlines with double newline
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(
    text: str,
    chunk_size: int = 1000,
    overlap: int = 200,
) -> list[ChunkSpec]:
    """
    Split text into overlapping chunks.

    Args:
        text: The text to chunk
        chunk_size: Maximum characters per chunk
        overlap: Number of characters to overlap between chunks

    Returns:
        List of ChunkSpec objects
    """
    text = normalize_whitespace(text)

    if not text:
        return []

    chunks = []
    start = 0
    index = 0

    while start < len(text):
        end = min(len(text), start + chunk_size)

        # Try to break at a sentence boundary if not at the end
        if end < len(text):
            # Look for sentence end within the last 20% of the chunk
            search_start = start + int(chunk_size * 0.8)
            for i in range(end, search_start, -1):
                if text[i - 1] in ".!?\n":
                    end = i
                    break

        chunk_text_content = text[start:end]

        chunks.append(
            ChunkSpec(
                index=index,
                text=chunk_text_content,
                char_start=start,
                char_end=end,
            )
        )

        index += 1

        # Move start forward, but ensure we make progress
        new_start = end - overlap
        if new_start <= start:
            new_start = end

        start = new_start

        # Safety check to prevent infinite loops
        if start >= len(text):
            break

    return chunks


# ============================================================================
# Type-aware chunking (v0.1)
# ============================================================================

# Target chunk size in characters (~300-600 tokens ≈ 1200-2400 chars)
TYPE_AWARE_CHUNK_SIZE = 1800
TYPE_AWARE_OVERLAP = 300

# Patterns for detecting section boundaries
CONTRACT_SECTION_PATTERNS = [
    # Numbered clauses: "1.", "1.1", "1.1.1", "Article 1", "Section 1"
    re.compile(
        r"^(?:article|section|clause)?\s*(\d+(?:\.\d+)*\.?)\s+(.+)", re.IGNORECASE | re.MULTILINE
    ),
    # Roman numerals: "I.", "II.", "III."
    re.compile(r"^((?:X{0,3})?(?:IX|IV|V?I{0,3})\.)\s+(.+)", re.MULTILINE),
    # Lettered sections: "(a)", "(b)", "(i)", "(ii)"
    re.compile(r"^\(([a-z]|[ivx]+)\)\s+(.+)", re.IGNORECASE | re.MULTILINE),
]

RFC_SECTION_PATTERNS = [
    # Markdown headings: "# Title", "## Section", "### Subsection"
    re.compile(r"^(#{1,6})\s+(.+)", re.MULTILINE),
    # Numbered sections: "1. Introduction", "1.1 Background"
    re.compile(r"^(\d+(?:\.\d+)*\.?)\s+([A-Z].+)", re.MULTILINE),
    # ALL CAPS headings
    re.compile(r"^([A-Z][A-Z\s]{4,})$", re.MULTILINE),
]

POLICY_SECTION_PATTERNS = [
    # Numbered sections
    re.compile(r"^(\d+(?:\.\d+)*\.?)\s+(.+)", re.MULTILINE),
    # Markdown headings
    re.compile(r"^(#{1,6})\s+(.+)", re.MULTILINE),
    # ALL CAPS headings
    re.compile(r"^([A-Z][A-Z\s]{4,})$", re.MULTILINE),
]


class Section:
    """Represents a detected section in a document."""

    def __init__(
        self,
        heading: str,
        level: int,
        start_pos: int,
        end_pos: int | None = None,
        content: str = "",
        number: str | None = None,
    ):
        self.heading = heading
        self.level = level
        self.start_pos = start_pos
        self.end_pos = end_pos
        self.content = content
        self.number = number
        self.children: list[Section] = []
        self.parent: Section | None = None

    def get_path(self) -> list[str]:
        """Get hierarchical path from root to this section."""
        path = []
        current = self
        while current:
            label = current.number or current.heading
            if label:
                path.insert(0, label.strip())
            current = current.parent
        return path


def detect_sections(text: str, patterns: list[re.Pattern]) -> list[Section]:
    """
    Detect section boundaries in text using the provided patterns.

    Returns a list of Section objects with detected boundaries.
    """
    matches = []

    for pattern in patterns:
        for match in pattern.finditer(text):
            groups = match.groups()
            if len(groups) >= 2:
                # Pattern has number/marker and title
                marker, title = groups[0], groups[1]
                heading = title.strip()
                number = marker.strip()
                # Determine level from marker
                if marker.startswith("#"):
                    level = len(marker)
                elif "." in marker:
                    level = marker.count(".") + 1
                else:
                    level = 1
            else:
                # Pattern has only heading (like ALL CAPS)
                heading = groups[0].strip()
                number = None
                level = 1

            matches.append(
                {
                    "heading": heading,
                    "number": number,
                    "level": level,
                    "start": match.start(),
                    "end": match.end(),
                }
            )

    # Sort by position and remove duplicates/overlaps
    matches.sort(key=lambda m: m["start"])
    filtered = []
    last_end = -1
    for m in matches:
        if m["start"] >= last_end:
            filtered.append(m)
            last_end = m["end"]

    # Create Section objects with end positions
    sections = []
    for i, m in enumerate(filtered):
        next_start = filtered[i + 1]["start"] if i + 1 < len(filtered) else len(text)
        section = Section(
            heading=m["heading"],
            level=m["level"],
            start_pos=m["start"],
            end_pos=next_start,
            content=text[m["end"] : next_start].strip(),
            number=m["number"],
        )
        sections.append(section)

    # Build hierarchy based on levels
    if sections:
        stack: list[Section] = []
        for section in sections:
            while stack and stack[-1].level >= section.level:
                stack.pop()
            if stack:
                section.parent = stack[-1]
                stack[-1].children.append(section)
            stack.append(section)

    return sections


def chunk_section(
    section: Section,
    chunk_size: int = TYPE_AWARE_CHUNK_SIZE,
    overlap: int = TYPE_AWARE_OVERLAP,
) -> list[ChunkSpec]:
    """
    Chunk a single section, preserving section_path metadata.

    If section content is small enough, returns a single chunk.
    Otherwise, splits into overlapping chunks with the same section_path.
    """
    text = section.content
    section_path = section.get_path()

    if not text.strip():
        return []

    # If content fits in one chunk, return as single chunk
    if len(text) <= chunk_size:
        return [
            ChunkSpec(
                index=0,
                text=text,
                char_start=section.start_pos,
                char_end=section.end_pos or (section.start_pos + len(text)),
                section_path=section_path,
            )
        ]

    # Otherwise, split into multiple chunks
    chunks = []
    start = 0
    idx = 0

    while start < len(text):
        end = min(len(text), start + chunk_size)

        # Try to break at sentence boundary
        if end < len(text):
            search_start = start + int(chunk_size * 0.7)
            for i in range(end, search_start, -1):
                if text[i - 1] in ".!?\n":
                    end = i
                    break

        chunk_content = text[start:end].strip()
        if chunk_content:
            chunks.append(
                ChunkSpec(
                    index=idx,
                    text=chunk_content,
                    char_start=section.start_pos + start,
                    char_end=section.start_pos + end,
                    section_path=section_path,
                )
            )
            idx += 1

        new_start = end - overlap
        if new_start <= start:
            new_start = end
        start = new_start

        if start >= len(text):
            break

    return chunks


def chunk_text_type_aware(
    text: str,
    doc_type: str | None = None,
    chunk_size: int = TYPE_AWARE_CHUNK_SIZE,
    overlap: int = TYPE_AWARE_OVERLAP,
) -> list[ChunkSpec]:
    """
    Split text into chunks based on document type.

    For CONTRACT, POLICY, and RFC documents, attempts to preserve
    semantic boundaries (clauses, sections, headings) and includes
    section_path metadata.

    For OTHER or unrecognized types, falls back to standard chunking.

    Args:
        text: The text to chunk
        doc_type: Document type (CONTRACT, POLICY, RFC, OTHER)
        chunk_size: Maximum characters per chunk
        overlap: Overlap between chunks

    Returns:
        List of ChunkSpec objects with section_path when available
    """
    text = normalize_whitespace(text)

    if not text:
        return []

    # Select patterns based on doc type
    patterns = []
    if doc_type == "CONTRACT":
        patterns = CONTRACT_SECTION_PATTERNS
    elif doc_type == "RFC":
        patterns = RFC_SECTION_PATTERNS
    elif doc_type == "POLICY":
        patterns = POLICY_SECTION_PATTERNS

    # Try type-aware chunking if we have patterns
    if patterns:
        sections = detect_sections(text, patterns)

        if sections:
            all_chunks = []
            global_idx = 0

            for section in sections:
                section_chunks = chunk_section(section, chunk_size, overlap)
                for chunk in section_chunks:
                    chunk.index = global_idx
                    all_chunks.append(chunk)
                    global_idx += 1

            # If we got chunks, return them
            if all_chunks:
                logger.info(
                    f"Type-aware chunking ({doc_type}): "
                    f"{len(sections)} sections, {len(all_chunks)} chunks"
                )
                return all_chunks

    # Fall back to standard chunking
    logger.info(f"Falling back to standard chunking for doc_type={doc_type}")
    return chunk_text(text, chunk_size, overlap)
